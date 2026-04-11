"""Office/PDF/HWPX 추출기 통합 테스트.

샘플 문서를 런타임에 생성한 뒤 각 추출기로 텍스트를 뽑고 기대 키워드가
포함되는지 검사한다. 실제 사용자 HWP 파일은 별도 fixture로 검증한다.
"""
from __future__ import annotations

from pathlib import Path

import pytest

from jean_finder.extractors.registry import build_default_registry

NEEDLE = "제안탐색기검증키워드"


@pytest.fixture(scope="module")
def registry():
    return build_default_registry(enable_ocr=False)


def test_docx_extraction(tmp_path_factory, registry):
    from docx import Document

    path = tmp_path_factory.mktemp("docx") / "sample.docx"
    doc = Document()
    doc.add_paragraph(f"머리말 {NEEDLE} 본문")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "셀A"
    table.rows[0].cells[1].text = "셀B"
    doc.save(str(path))

    result = registry.for_path(path).extract(path)
    assert NEEDLE in result.text
    assert "셀A" in result.text and "셀B" in result.text


def test_xlsx_extraction(tmp_path_factory, registry):
    from openpyxl import Workbook

    path = tmp_path_factory.mktemp("xlsx") / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "데이터"
    ws.append(["헤더1", "헤더2"])
    ws.append([NEEDLE, 42])
    wb.save(str(path))

    result = registry.for_path(path).extract(path)
    assert NEEDLE in result.text
    assert "데이터" in result.text


def test_pptx_extraction(tmp_path_factory, registry):
    from pptx import Presentation
    from pptx.util import Inches

    path = tmp_path_factory.mktemp("pptx") / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"제목 {NEEDLE}"
    prs.save(str(path))

    result = registry.for_path(path).extract(path)
    assert NEEDLE in result.text
    assert result.page_count == 1


def test_pdf_extraction(tmp_path_factory, registry):
    # PyMuPDF의 insert_text 기본 폰트(Helvetica)는 한글 글리프가 없어
    # 비ASCII 문자가 깨지므로 PDF 검증에는 ASCII needle을 사용한다.
    import fitz

    ascii_needle = "JeAnFinderPdfNeedle"
    path = tmp_path_factory.mktemp("pdf") / "sample.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), f"hello {ascii_needle} world")
    doc.save(str(path))
    doc.close()

    result = registry.for_path(path).extract(path)
    assert ascii_needle in result.text
    assert result.page_count == 1
    assert result.ocr_used is False


def test_hwpx_extraction(tmp_path_factory, registry):
    """HWPX는 단순 ZIP+XML이므로 최소 구조의 가짜 hwpx를 만들어 검증한다."""
    import zipfile

    path = tmp_path_factory.mktemp("hwpx") / "sample.hwpx"
    section_xml = (
        "<?xml version='1.0' encoding='UTF-8'?>"
        "<hp:sec xmlns:hp='http://www.hancom.co.kr/hwpml/2011/paragraph'>"
        f"<hp:p><hp:run><hp:t>{NEEDLE}</hp:t></hp:run></hp:p>"
        "</hp:sec>"
    )
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("Contents/section0.xml", section_xml)

    result = registry.for_path(path).extract(path)
    assert NEEDLE in result.text


def test_registry_unknown_extension(registry):
    assert registry.for_path(Path("foo.txt")) is None

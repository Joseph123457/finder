from __future__ import annotations

from pathlib import Path

from .base import BaseExtractor, ExtractionResult, ExtractorError


class PptxExtractor(BaseExtractor):
    name = "pptx"
    extensions = (".pptx",)

    def extract(self, path: Path) -> ExtractionResult:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ExtractorError(f"python-pptx not installed: {e}")

        try:
            prs = Presentation(str(path))
        except Exception as e:
            raise ExtractorError(f"PPTX open failed: {e}")

        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            parts.append(f"# Slide {idx}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text:
                            parts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text]
                        if row_cells:
                            parts.append("\t".join(row_cells))
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            if notes:
                parts.append(f"[Notes] {notes}")

        return ExtractionResult(text="\n".join(parts), page_count=len(prs.slides))
